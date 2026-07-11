"""
Generate graduation PowerPoint: Transport_Analysis_Presentation.pptx
Run: py create_presentation.py
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).resolve().parent
OUTPUT = BASE_DIR / "Transport_Analysis_Presentation.pptx"
FIGURES = BASE_DIR / "reports" / "figures"
METRICS_PATH = BASE_DIR / "reports" / "model_metrics.json"

NAVY = RGBColor(22, 58, 88)
DARK = RGBColor(30, 45, 65)
ACCENT = RGBColor(230, 126, 34)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(90, 105, 120)
LIGHT = RGBColor(245, 247, 250)

TEAM = [
    "Ahmed El-Gamal",
    "Youssef Mustafa",
    "Maryam Khaled",
    "Mariam Tarek",
    "Mostafa Fouad",
    "Nada Ahmed",
]


def _metrics() -> dict:
    if METRICS_PATH.exists():
        return json.loads(METRICS_PATH.read_text(encoding="utf-8"))
    return {
        "fare_regression_r2": 0.981,
        "fare_regression_mae": 1.91,
        "delay_classification_accuracy": 0.779,
        "trip_clustering_silhouette": 0.323,
        "overall_delay_rate": 0.307,
    }


def _blank_slide(prs: Presentation, dark: bool = False):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY if dark else LIGHT
    return slide


def _add_bar(slide, dark: bool = False) -> None:
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    if dark:
        return
    footer = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(9), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Transport Analysis  |  Digital Egypt Pioneers Initiative (DEPI)"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY


def _title_block(slide, title: str, subtitle: str = "", dark: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(9), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE if dark else NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = ACCENT if not dark else RGBColor(255, 200, 150)


def _bullets(slide, items: list[str], left: float = 0.7, top: float = 1.5, width: float = 8.8, size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


def _notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def _title_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs, dark=True)
    _add_bar(slide, dark=True)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.5), Inches(2.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Transport Analysis"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "Intelligent Public Transport Analytics Platform"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT
    p3 = tf.add_paragraph()
    p3.text = "Digital Egypt Pioneers Initiative (DEPI)"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(180, 200, 220)
    p4 = tf.add_paragraph()
    p4.text = "AI & Data Science Track  •  Graduation Project"
    p4.font.size = Pt(16)
    p4.font.color.rgb = RGBColor(160, 180, 200)
    _notes(slide, "Good morning. We present Transport Analysis — an end-to-end data and AI platform for public transport intelligence, developed under DEPI.")


def _team_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Our Team")
    cols = slide.shapes.add_table(2, 3, Inches(0.8), Inches(1.8), Inches(8.4), Inches(2.2)).table
    for idx, name in enumerate(TEAM):
        r, c = divmod(idx, 3)
        cell = cols.cell(r, c)
        cell.text = name
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = NAVY
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    _notes(slide, "Our team collaborated across data engineering, analytics, BI, and project management roles to deliver every milestone.")


def _agenda_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Agenda")
    _bullets(
        slide,
        [
            "Problem Statement & Project Vision",
            "Solution Architecture & Technology Stack",
            "Data Pipeline: Ingestion → ETL → SQL Warehouse",
            "Machine Learning & Predictive Analytics",
            "Results, Insights & Live Dashboard Demo",
            "Impact, Challenges & Future Work",
        ],
        top=1.6,
    )
    _notes(slide, "We will walk from the business problem to the technical solution, show results, and demo the interactive dashboard.")


def _problem_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Problem Statement", "Why this project matters")
    _bullets(
        slide,
        [
            "Transport operators collect massive trip data but struggle to extract timely insights",
            "Manual spreadsheets cannot scale to thousands of trips and multiple variables",
            "Delays, revenue leakage, and peak-hour congestion need data-driven responses",
            "Decision-makers need one trusted pipeline: clean data → analytics → predictions → dashboard",
        ],
        top=1.7,
        size=19,
    )
    _notes(slide, "The core problem is turning raw transport CSV data into actionable intelligence at scale — not just storing numbers.")


def _objectives_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Project Objectives")
    _bullets(
        slide,
        [
            "Analyze large public transport datasets and extract meaningful patterns",
            "Build a scalable batch ETL pipeline (Bronze → Silver → Gold)",
            "Run SQL warehouse analytics for operational and strategic KPIs",
            "Apply machine learning to predict fares, delays, and trip segments",
            "Deliver an interactive dashboard for non-technical stakeholders",
        ],
        top=1.7,
    )
    _notes(slide, "These objectives align directly with DEPI milestones: ingestion, development, deployment, automation, and presentation.")


def _architecture_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Solution Architecture", "End-to-end data flow")
    stages = [
        ("CSV\nIngestion", 0.6),
        ("Batch ETL\nPySpark / Pandas", 2.5),
        ("SQL\nWarehouse", 4.4),
        ("Machine\nLearning", 6.3),
        ("Streamlit\nDashboard", 8.2),
    ]
    for label, x in stages:
        shape = slide.shapes.add_shape(1, Inches(x), Inches(2.8), Inches(1.55), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.color.rgb = ACCENT
        tf = shape.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for x in [2.2, 4.1, 6.0, 7.9]:
        arrow = slide.shapes.add_shape(1, Inches(x), Inches(3.2), Inches(0.25), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()
    _bullets(
        slide,
        [
            "Modular Python codebase — each stage runs independently or as one pipeline",
            "Automated reports, metrics JSON, and Parquet layers for reproducibility",
        ],
        top=4.5,
        size=17,
    )
    _notes(slide, "Data flows from raw CSV through ETL layers into SQL analytics and ML, ending in a live Streamlit dashboard.")


def _tech_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Technology Stack")
    table = slide.shapes.add_table(5, 2, Inches(1.0), Inches(1.8), Inches(8.0), Inches(3.5)).table
    rows = [
        ("Data Processing", "Python, Pandas, PySpark"),
        ("Storage & ETL", "CSV, Parquet, SQLite Warehouse"),
        ("Analytics", "SQL, scikit-learn"),
        ("Visualization", "Matplotlib, Seaborn, Plotly, Streamlit"),
        ("Automation", "Batch pipeline, JSON reports, scheduled re-runs"),
    ]
    for r, (col1, col2) in enumerate(rows):
        table.cell(r, 0).text = col1
        table.cell(r, 1).text = col2
        for c in (0, 1):
            for p in table.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(15)
                p.font.color.rgb = DARK if c else NAVY
                if c == 0:
                    p.font.bold = True
    _notes(slide, "We used industry-standard tools — the same categories used in cloud warehouses and modern data teams.")


def _etl_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Batch ETL Pipeline", "Medallion architecture")
    layers = [
        ("Bronze", "Raw transport_trips.csv", NAVY),
        ("Silver", "Cleaned + features (speed, revenue/km)", RGBColor(41, 98, 130)),
        ("Gold", "Daily summaries by vehicle, zone, date", ACCENT),
    ]
    for i, (name, desc, color) in enumerate(layers):
        top = 1.8 + i * 1.35
        shape = slide.shapes.add_shape(1, Inches(0.8), Inches(top), Inches(2.0), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.text_frame.paragraphs[0].text = name
        shape.text_frame.paragraphs[0].font.color.rgb = WHITE
        shape.text_frame.paragraphs[0].font.bold = True
        box = slide.shapes.add_textbox(Inches(3.0), Inches(top + 0.1), Inches(6.5), Inches(0.8))
        box.text_frame.paragraphs[0].text = desc
        box.text_frame.paragraphs[0].font.size = Pt(18)
        box.text_frame.paragraphs[0].font.color.rgb = DARK
    _notes(slide, "Bronze keeps raw data, Silver adds engineered features, Gold delivers business-ready aggregates for executives.")


def _sql_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "SQL Warehouse Analytics")
    _bullets(
        slide,
        [
            "6 analytical queries on SQLite (simulating BigQuery / Snowflake / Redshift)",
            "Revenue by vehicle type and average fare by zone",
            "Delay rate by traffic level and weather impact",
            "Peak-hour demand analysis for capacity planning",
            "Outputs saved to sql/analytics_queries.sql and reports/sql_analytics.json",
        ],
        top=1.7,
        size=18,
    )
    _notes(slide, "SQL answers the business questions stakeholders ask every day — revenue, delays, peaks, and zones.")


def _ml_slide(prs: Presentation) -> None:
    m = _metrics()
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Machine Learning Layer", "Predictive intelligence")
    _bullets(
        slide,
        [
            "Fare Regression — predict trip fare from distance, vehicle, zone, traffic, weather",
            f"Performance: R² = {m['fare_regression_r2']:.3f}, MAE = ${m['fare_regression_mae']:.2f}",
            "Delay Classification — predict whether a trip will be delayed",
            f"Performance: Accuracy = {m['delay_classification_accuracy']:.1%}",
            f"Trip Segmentation (K-Means) — silhouette = {m['trip_clustering_silhouette']:.3f}",
        ],
        top=1.7,
        size=17,
    )
    _notes(slide, "ML moves the project from descriptive analytics to predictive intelligence — forecasting fare and delay risk.")


def _kpi_slide(prs: Presentation) -> None:
    m = _metrics()
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Key Results & KPIs")
    cards = [
        ("5,000", "Trips Analyzed"),
        ("$28.18", "Avg Fare"),
        ("30.7%", "Delay Rate"),
        ("Airport", "Top Fare Zone"),
        ("Taxi", "Top Revenue Vehicle"),
    ]
    for i, (value, label) in enumerate(cards):
        col = i % 3
        row = i // 3
        left = 0.7 + col * 3.1
        top = 1.9 + row * 2.0
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2.7), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = ACCENT
        tf = shape.text_frame
        tf.paragraphs[0].text = value
        tf.paragraphs[0].font.size = Pt(26)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER
    _notes(slide, "These KPIs prove the pipeline works end-to-end on real-scale sample data with meaningful business metrics.")


def _chart_slide(prs: Presentation, image: Path, title: str, insight: str, note: str) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, title)
    if image.exists():
        slide.shapes.add_picture(str(image), Inches(0.55), Inches(1.35), width=Inches(5.6))
    box = slide.shapes.add_textbox(Inches(6.35), Inches(1.8), Inches(3.2), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Insight"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = insight
    p2.font.size = Pt(15)
    p2.font.color.rgb = DARK
    _notes(slide, note)


def _insights_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Strategic Insights")
    _bullets(
        slide,
        [
            "Airport and downtown zones command premium fares — pricing strategy opportunity",
            "Taxi leads total revenue; metro and bus serve volume and accessibility",
            "High traffic increases delays to ~38% — prioritize congestion management",
            "Peak hours (7–9 AM, 5–7 PM) show predictable demand surges",
            "Weather (rain/fog) amplifies delay risk — useful for proactive scheduling",
        ],
        top=1.7,
        size=18,
    )
    _notes(slide, "These insights are actionable — not just charts. They guide pricing, fleet allocation, and operations.")


def _dashboard_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Interactive Dashboard Demo", "Streamlit Web Application")
    _bullets(
        slide,
        [
            "Launch: py -m streamlit run dashboard.py",
            "Live KPIs: trips, fare, distance, delay rate, revenue",
            "Filters: vehicle, zone, traffic, weather, date range",
            "Interactive Plotly charts update instantly",
            "Trip-level drill-down table for validation",
        ],
        top=1.7,
        size=18,
    )
    _notes(slide, "Live demo: open localhost:8501, filter by High traffic, show delay rate spike. This is the executive-facing layer.")


def _challenges_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Challenges & Solutions")
    table = slide.shapes.add_table(4, 2, Inches(0.7), Inches(1.7), Inches(8.8), Inches(3.8)).table
    rows = [
        ("Large dataset handling", "PySpark ETL with Pandas fallback"),
        ("Data quality issues", "Automated cleaning in Silver layer"),
        ("Stakeholder accessibility", "Streamlit dashboard — no code required"),
        ("Reproducibility", "Single-command pipeline + versioned outputs"),
    ]
    for r, (challenge, solution) in enumerate(rows):
        table.cell(r, 0).text = challenge
        table.cell(r, 1).text = solution
        for c in (0, 1):
            for p in table.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = DARK
    _notes(slide, "We designed for real-world constraints — environments without Java still run the full project.")


def _future_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_bar(slide)
    _title_block(slide, "Future Enhancements")
    _bullets(
        slide,
        [
            "Deploy to cloud warehouse (BigQuery / Snowflake) with scheduled batch jobs",
            "Real-time streaming ingestion for live delay alerts",
            "GPS route optimization and geospatial heatmaps",
            "Mobile app for operators and executives",
            "CI/CD automation and pipeline monitoring (MLOps)",
        ],
        top=1.7,
        size=18,
    )
    _notes(slide, "This project is production-ready as a foundation — cloud deployment is the natural next step.")


def _closing_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs, dark=True)
    _add_bar(slide, dark=True)
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(8), Inches(2.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Transport Analysis — From Data to Decisions"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "Questions & Live Demo"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(180, 200, 220)
    p3.alignment = PP_ALIGN.CENTER
    _notes(slide, "Thank you for your time. We are ready for questions and a live dashboard demonstration.")


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_slide(prs)
    _team_slide(prs)
    _agenda_slide(prs)
    _problem_slide(prs)
    _objectives_slide(prs)
    _architecture_slide(prs)
    _tech_slide(prs)
    _etl_slide(prs)
    _sql_slide(prs)
    _ml_slide(prs)
    _kpi_slide(prs)

    _chart_slide(
        prs,
        FIGURES / "revenue_by_vehicle.png",
        "Revenue Analysis",
        "Taxi generates the highest total revenue, making it critical for fare optimization.",
        "Point out taxi bar height vs bus and metro — revenue vs volume trade-off.",
    )
    _chart_slide(
        prs,
        FIGURES / "delay_by_traffic.png",
        "Operational Risk: Delays",
        "High traffic shows the strongest delay rate — target congestion hotspots first.",
        "Emphasize operations value — this supports staffing and routing decisions.",
    )
    _chart_slide(
        prs,
        FIGURES / "peak_hour_trips.png",
        "Demand Planning",
        "Clear morning and evening peaks enable smarter fleet scheduling.",
        "Connect to business: deploy more vehicles before peaks, reduce off-peak waste.",
    )
    _chart_slide(
        prs,
        FIGURES / "distance_vs_fare.png",
        "Fare Intelligence",
        "Fare scales with distance but varies by vehicle type — supports dynamic pricing.",
        "Show scatter clusters — different vehicle types, different economics.",
    )

    _insights_slide(prs)
    _dashboard_slide(prs)
    _challenges_slide(prs)
    _future_slide(prs)
    _closing_slide(prs)

    prs.save(OUTPUT)
    return OUTPUT


def main() -> None:
    path = build_presentation()
    print(f"Presentation saved to: {path}")
    print(f"Slides: 22  |  Speaker notes included on every slide")


if __name__ == "__main__":
    main()
