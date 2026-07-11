"""
Generate a light, professional PowerPoint for Transport Analysis.
Includes UML diagrams (component, class, sequence).
Output: Transport_Analysis_Presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from uml_diagrams import generate_all as generate_uml_diagrams
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt, Emu

BASE = Path(__file__).resolve().parent
FIG = BASE / "reports" / "figures"
UML = BASE / "reports" / "uml"
PHOTO = BASE / "photo"
OUT = BASE / "Transport_Analysis_Presentation.pptx"
OUT_TMP = BASE / "Transport_Analysis_Presentation_new.pptx"
TOTAL_SLIDES = 33

# Light professional palette
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
SLATE = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xFA, 0xFC)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT_SOFT = RGBColor(0xCC, 0xFB, 0xF1)
GOLD = RGBColor(0xB4, 0x53, 0x09)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape) -> None:
    shape.line.fill.background()


def set_run(run, text: str, size: int, bold: bool = False, color: RGBColor = SLATE, font: str = "Calibri") -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    run = p.add_run()
    set_run(run, text, size, bold, color, font)
    return box


def text_frame_set(tf, lines: list[tuple[str, int, bool, RGBColor]], align=PP_ALIGN.LEFT, space_after=6):
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.clear()
        run = p.add_run()
        set_run(run, text, size, bold, color)


def footer(slide, page: int, total: int = TOTAL_SLIDES) -> None:
    # bottom bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), SLIDE_W, Inches(0.35))
    rgb_fill(bar, BG)
    no_line(bar)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.12), Inches(12.333), Pt(1.5))
    rgb_fill(line, LINE)
    no_line(line)
    add_textbox(slide, Inches(0.5), Inches(7.18), Inches(8), Inches(0.28),
                "Transport Analysis  ·  DEPI AI & Data Science", 10, False, MUTED)
    add_textbox(slide, Inches(11.2), Inches(7.18), Inches(1.6), Inches(0.28),
                f"{page} / {total}", 10, False, MUTED, PP_ALIGN.RIGHT)


def top_accent(slide) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
    rgb_fill(bar, TEAL)
    no_line(bar)


def bg_slide(slide) -> None:
    # white/light background via full rect
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rgb_fill(shape, WHITE)
    no_line(shape)
    # send to back by reordering is hard; draw first so content sits on top
    top_accent(slide)


def card(slide, left, top, width, height, fill=CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rgb_fill(shape, fill)
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    # slight corner
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def accent_bar_left(slide, left, top, height, color=TEAL):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    rgb_fill(bar, color)
    no_line(bar)


def bullet_block(slide, left, top, width, height, items: list[str], size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        p.clear()
        run = p.add_run()
        set_run(run, f"•  {item}", size, False, SLATE)
    return box


def section_slide(prs, part: str, title: str, subtitle: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(slide)
    # soft teal panel
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), SLIDE_W, Inches(3.0))
    rgb_fill(panel, RGBColor(0xF0, 0xFD, 0xFA))
    no_line(panel)
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(0.4),
                part, 14, True, TEAL)
    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.8),
                title, 40, True, NAVY)
    add_textbox(slide, Inches(0.8), Inches(3.9), Inches(11), Inches(0.5),
                subtitle, 18, False, MUTED)
    footer(slide, page)
    return slide


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(slide)
    # left accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), SLIDE_H)
    rgb_fill(strip, TEAL)
    no_line(strip)
    soft = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), 0, Inches(0.08), SLIDE_H)
    rgb_fill(soft, ACCENT_SOFT)
    no_line(soft)

    add_textbox(slide, Inches(0.9), Inches(1.8), Inches(11), Inches(0.4),
                "DIGITAL EGYPT PIONEERS INITIATIVE (DEPI)", 13, True, TEAL)
    add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.0),
                "Transport Analysis", 48, True, NAVY)
    add_textbox(slide, Inches(0.9), Inches(3.4), Inches(11), Inches(0.5),
                "End-to-End Public Transport Analytics Pipeline", 22, False, SLATE)
    add_textbox(slide, Inches(0.9), Inches(4.2), Inches(11), Inches(0.4),
                "AI & Data Science Track  ·  Graduation Project", 16, False, MUTED)
    add_textbox(slide, Inches(0.9), Inches(5.2), Inches(11), Inches(0.4),
                "Python  ·  Pandas  ·  PySpark  ·  SQL  ·  scikit-learn  ·  Streamlit", 14, False, TEAL)
    footer(slide, 1)
    return slide


def content_title(slide, title: str, subtitle: str | None = None):
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(12), Inches(0.55),
                title, 28, True, NAVY)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.9), Inches(12), Inches(0.35),
                    subtitle, 13, False, MUTED)


def notes(slide, text: str):
    ns = slide.notes_slide
    ns.notes_text_frame.text = text


def uml_slide(prs, title: str, subtitle: str, image_path: Path, page: int, note: str = "") -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, title, subtitle)
    if image_path.exists():
        s.shapes.add_picture(str(image_path), Inches(0.45), Inches(1.25), width=Inches(12.4))
    footer(s, page)
    if note:
        notes(s, note)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    uml_paths = generate_uml_diagrams()

    # 1 Title
    title_slide(prs)

    # 2 Agenda
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Agenda")
    items = [
        ("01", "Problem & objectives"),
        ("02", "Solution architecture, UML diagrams & technology stack"),
        ("03", "Data design, ETL, and SQL analytics"),
        ("04", "Machine learning models & metrics"),
        ("05", "Reporting, dashboard & live demo"),
        ("06", "Insights, impact & next steps"),
    ]
    for i, (num, label) in enumerate(items):
        y = Inches(1.4) + Inches(i * 0.8)
        card(s, Inches(0.7), y, Inches(11.9), Inches(0.68))
        add_textbox(s, Inches(1.0), y + Inches(0.15), Inches(0.8), Inches(0.4), num, 18, True, TEAL)
        add_textbox(s, Inches(2.0), y + Inches(0.18), Inches(10), Inches(0.4), label, 18, False, SLATE)
    footer(s, 2)
    notes(s, "Walk through the agenda so the audience knows the structure of the talk.")

    # 3 Team
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Project Team", "DEPI — AI & Data Science Track")
    team = [
        ("Ahmed El-Gamal", "Pipeline / ETL"),
        ("Youssef Mustafa", "SQL Analytics"),
        ("Maryam Khaled", "Machine Learning"),
        ("Mariam Tarek", "Reporting"),
        ("Mostafa Fouad", "Dashboard / UI"),
        ("Nada Ahmed", "Documentation"),
    ]
    for i, (name, role) in enumerate(team):
        col, row = i % 3, i // 3
        left = Inches(0.7) + Inches(col * 4.1)
        top = Inches(1.6) + Inches(row * 2.3)
        card(s, left, top, Inches(3.85), Inches(1.9))
        accent_bar_left(s, left, top, Inches(1.9))
        add_textbox(s, left + Inches(0.3), top + Inches(0.5), Inches(3.3), Inches(0.5), name, 16, True, NAVY)
        add_textbox(s, left + Inches(0.3), top + Inches(1.05), Inches(3.3), Inches(0.4), role, 13, False, TEAL)
    footer(s, 3)
    notes(s, "Briefly introduce each team member and their focus area.")

    # 4 Part 1
    section_slide(prs, "PART 1", "The Problem", "Why transport analytics matters", 4)

    # 5 Challenge
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "The Challenge")
    challenges = [
        "Public transport operators generate large volumes of trip data daily",
        "Fares, delays, zones, weather, and traffic interact in complex ways",
        "Manual reporting is slow and does not support predictive decisions",
        "Need one pipeline: clean data → SQL insights → ML forecasts → interactive UI",
    ]
    for i, c in enumerate(challenges):
        y = Inches(1.4) + Inches(i * 1.15)
        card(s, Inches(0.7), y, Inches(11.9), Inches(0.95))
        num_box = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.22), Inches(0.5), Inches(0.5))
        rgb_fill(num_box, TEAL)
        no_line(num_box)
        add_textbox(s, Inches(1.0), y + Inches(0.3), Inches(0.5), Inches(0.4), str(i + 1), 14, True, WHITE, PP_ALIGN.CENTER)
        add_textbox(s, Inches(1.8), y + Inches(0.3), Inches(10.3), Inches(0.5), c, 16, False, SLATE)
    footer(s, 5)
    notes(s, "Frame the business problem: operators need faster, predictive insight from trip data.")

    # 6 Objectives
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Project Objectives")
    objs = [
        ("Ingest", "Load or generate realistic trip datasets"),
        ("Transform", "Bronze → Silver → Gold medallion ETL"),
        ("Analyze", "SQL warehouse queries for ops KPIs"),
        ("Predict", "Fare, delay, and trip segmentation models"),
        ("Report", "Automated Markdown + PNG charts"),
        ("Explore", "Interactive Streamlit dashboard"),
    ]
    for i, (t, d) in enumerate(objs):
        col, row = i % 3, i // 3
        left = Inches(0.6) + Inches(col * 4.15)
        top = Inches(1.5) + Inches(row * 2.4)
        card(s, left, top, Inches(3.95), Inches(2.1))
        add_textbox(s, left + Inches(0.25), top + Inches(0.4), Inches(3.4), Inches(0.45), t, 18, True, TEAL)
        add_textbox(s, left + Inches(0.25), top + Inches(1.0), Inches(3.4), Inches(0.7), d, 14, False, SLATE)
    footer(s, 6)

    # 7 Part 2
    section_slide(prs, "PART 2", "Architecture", "How the system is designed", 7)

    # 8 Architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "End-to-End Architecture")
    stages = [
        ("1. Ingestion", "data_pipeline.py"),
        ("2. Batch ETL", "spark_etl.py"),
        ("3. SQL", "sql_analytics.py"),
        ("4. ML", "train_ai_model.py"),
        ("5. Report", "analysis_report.py"),
        ("6. UI", "dashboard.py"),
    ]
    for i, (t, m) in enumerate(stages):
        left = Inches(0.45) + Inches(i * 2.15)
        card(s, left, Inches(1.5), Inches(2.0), Inches(2.0))
        add_textbox(s, left + Inches(0.1), Inches(1.85), Inches(1.8), Inches(0.5), t, 13, True, NAVY, PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.1), Inches(2.55), Inches(1.8), Inches(0.5), m, 11, False, TEAL, PP_ALIGN.CENTER)
        if i < 5:
            add_textbox(s, left + Inches(1.85), Inches(2.2), Inches(0.35), Inches(0.4), "→", 18, True, MUTED, PP_ALIGN.CENTER)
    bullet_block(s, Inches(0.7), Inches(4.0), Inches(11.5), Inches(2.5), [
        "Orchestrator: src/main.py runs stages 1–5 in order",
        "Dashboard is a separate entry point reusing the same dataset",
        "Outputs land in data/ and reports/ for demos and grading",
    ], 16)
    footer(s, 8)

    # 9–11 UML diagrams
    uml_slide(
        prs,
        "UML Component Diagram",
        "Modules, data stores, and how main.py orchestrates the pipeline",
        uml_paths["component"],
        9,
        "Explain each Python module as a component and how data flows between layers.",
    )
    uml_slide(
        prs,
        "UML Class Diagram",
        "Core classes and public functions in each source module",
        uml_paths["class"],
        10,
        "Walk through DataPaths, pipeline functions, and how main.py uses each module.",
    )
    uml_slide(
        prs,
        "UML Sequence Diagram",
        "Step-by-step execution when running py src/main.py",
        uml_paths["sequence"],
        11,
        "Trace the call order: load data → ETL → SQL → ML → report → console summary.",
    )

    # 12 Tech stack
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Technology Stack")
    stacks = [
        ("Data", ["Python 3.10+", "Pandas / NumPy", "PyArrow (Parquet)"]),
        ("ETL", ["PySpark (optional)", "Pandas fallback", "Medallion layers"]),
        ("Analytics", ["SQLite warehouse", "SQL aggregations", "JSON exports"]),
        ("ML", ["scikit-learn", "RandomForest", "KMeans clustering"]),
        ("Viz / UI", ["Matplotlib / Seaborn", "Streamlit", "Plotly charts"]),
    ]
    for i, (title, lines) in enumerate(stacks):
        left = Inches(0.45) + Inches(i * 2.55)
        card(s, left, Inches(1.45), Inches(2.4), Inches(4.8))
        header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.45), Inches(2.4), Inches(0.7))
        rgb_fill(header, NAVY)
        no_line(header)
        add_textbox(s, left, Inches(1.6), Inches(2.4), Inches(0.45), title, 15, True, WHITE, PP_ALIGN.CENTER)
        for j, line in enumerate(lines):
            add_textbox(s, left + Inches(0.15), Inches(2.5) + Inches(j * 0.7), Inches(2.1), Inches(0.5), line, 13, False, SLATE, PP_ALIGN.CENTER)
    footer(s, 12)

    # 13 Repo structure
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Repository Structure")
    card(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(5.2))
    structure = [
        "src/              main, data_pipeline, spark_etl, sql_analytics, train_ai_model, analysis_report",
        "data/raw          transport_trips.csv (Bronze)",
        "data/silver       cleaned Parquet with engineered features",
        "data/gold         daily aggregated summaries",
        "data/warehouse    transport.db (SQLite analytics)",
        "sql/              analytics_queries.sql",
        "reports/          metrics JSON, clustered CSV, figures, Markdown report",
        "dashboard.py      Streamlit interactive website",
    ]
    for i, line in enumerate(structure):
        add_textbox(s, Inches(1.1), Inches(1.7) + Inches(i * 0.55), Inches(11), Inches(0.45), line, 14, False, SLATE, font="Consolas")
    footer(s, 13)

    # 14 Part 3
    section_slide(prs, "PART 3", "Data & ETL", "From raw trips to gold summaries", 14)

    # 15 Schema
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Trip Data Schema", "5,000 synthetic trips  ·  seed = 42 for reproducibility")
    card(s, Inches(0.55), Inches(1.45), Inches(6.0), Inches(5.1))
    add_textbox(s, Inches(0.85), Inches(1.7), Inches(5.4), Inches(0.4), "Core Fields", 16, True, TEAL)
    bullet_block(s, Inches(0.85), Inches(2.25), Inches(5.4), Inches(3.8), [
        "trip_id, trip_date, pickup_hour",
        "vehicle_type: Taxi / Bus / Metro / RideShare",
        "zone: Downtown, Airport, Suburb, Industrial, Coastal",
        "trip_distance_km, trip_duration_min",
        "passenger_count",
    ], 14)
    card(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(5.1))
    add_textbox(s, Inches(7.1), Inches(1.7), Inches(5.4), Inches(0.4), "Fare, Context & Targets", 16, True, TEAL)
    bullet_block(s, Inches(7.1), Inches(2.25), Inches(5.4), Inches(3.8), [
        "base_fare, tip_amount, total_fare (ML target)",
        "payment_type: Cash / Card / Mobile",
        "weather, traffic_level, delay_minutes",
        "is_delayed (1 if delay > 5 min) — ML target",
    ], 14)
    footer(s, 15)

    # 16 Synthetic design
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Synthetic Data Design Logic")
    rules = [
        ("Distance", "Gamma distribution per vehicle type, clipped to 0.5–45 km"),
        ("Duration", "distance × random factor × traffic × weather multipliers"),
        ("Fare", "distance × rate_per_km × zone multiplier + tip by payment type"),
        ("Delay", "traffic × Exp(2.5) + rain/fog bonus; delayed if > 5 minutes"),
        ("Reproducibility", "Fixed seed (42); load CSV if present, else generate"),
    ]
    for i, (t, d) in enumerate(rules):
        y = Inches(1.35) + Inches(i * 1.0)
        card(s, Inches(0.7), y, Inches(11.9), Inches(0.85))
        add_textbox(s, Inches(1.0), y + Inches(0.22), Inches(2.5), Inches(0.45), t, 15, True, TEAL)
        add_textbox(s, Inches(3.6), y + Inches(0.22), Inches(8.6), Inches(0.5), d, 14, False, SLATE)
    footer(s, 16)

    # 17 Medallion
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Medallion ETL: Bronze → Silver → Gold")
    layers = [
        ("BRONZE", "Raw CSV ingest", ["transport_trips.csv", "5,000 rows", "As-recorded trips"], RGBColor(0xB4, 0x53, 0x09)),
        ("SILVER", "Clean + features", ["dropna + clean", "speed_kmh", "revenue_per_km", "is_peak_hour", "Parquet output"], TEAL),
        ("GOLD", "Daily aggregates", ["groupby date × vehicle × zone", "trip_count, revenue", "avg_fare, delay_rate", "3,522 summary rows"], NAVY),
    ]
    for i, (name, sub, lines, color) in enumerate(layers):
        left = Inches(0.55) + Inches(i * 4.2)
        card(s, left, Inches(1.45), Inches(3.95), Inches(5.1))
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.45), Inches(3.95), Inches(1.1))
        rgb_fill(hdr, color)
        no_line(hdr)
        add_textbox(s, left, Inches(1.55), Inches(3.95), Inches(0.45), name, 18, True, WHITE, PP_ALIGN.CENTER)
        add_textbox(s, left, Inches(2.05), Inches(3.95), Inches(0.35), sub, 12, False, WHITE, PP_ALIGN.CENTER)
        for j, line in enumerate(lines):
            add_textbox(s, left + Inches(0.3), Inches(2.9) + Inches(j * 0.55), Inches(3.4), Inches(0.45), f"•  {line}", 13, False, SLATE)
    footer(s, 17)

    # 18 Feature eng
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Feature Engineering & Engine Fallback")
    bullet_block(s, Inches(0.7), Inches(1.5), Inches(11.5), Inches(5), [
        "speed_kmh = distance / (duration in hours)",
        "revenue_per_km = total_fare / distance",
        "is_peak_hour = 1 if pickup hour is 7–9 or 17–19",
        "run_batch_etl() probes Java → PySpark; otherwise identical Pandas path",
        "Writes reports/etl_report.json (engine, row counts, output paths)",
        "Latest run: engine = pandas · 5,000 silver · 3,522 gold rows",
    ], 17)
    footer(s, 18)

    # 19 Part 4
    section_slide(prs, "PART 4", "SQL Analytics", "Warehouse queries and business KPIs", 19)

    # 20 Warehouse
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "SQLite Analytics Warehouse")
    steps = [
        "Load Silver Parquet (or raw CSV fallback) into data/warehouse/transport.db",
        "Replace table transport_trips on each pipeline run",
        "Execute 6 analytical queries via pandas.read_sql_query",
        "Export results → reports/sql_analytics.json",
        "Regenerate sql/analytics_queries.sql for documentation and reuse",
    ]
    for i, step in enumerate(steps):
        y = Inches(1.4) + Inches(i * 0.95)
        card(s, Inches(0.7), y, Inches(11.9), Inches(0.8))
        add_textbox(s, Inches(1.0), y + Inches(0.2), Inches(0.6), Inches(0.4), f"{i+1:02d}", 16, True, TEAL)
        add_textbox(s, Inches(1.8), y + Inches(0.22), Inches(10.3), Inches(0.45), step, 15, False, SLATE)
    footer(s, 20)

    # 21 Six queries
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Six Analytical SQL Queries")
    queries = [
        ("daily_trip_volume", "Trips & revenue per day"),
        ("avg_fare_by_zone", "Premium zones by avg fare"),
        ("delay_rate_by_traffic", "Delay hotspot by traffic"),
        ("revenue_by_vehicle_type", "Revenue ranking by mode"),
        ("peak_hour_analysis", "Demand & delay by hour"),
        ("weather_impact", "Weather vs delay risk"),
    ]
    for i, (name, desc) in enumerate(queries):
        col, row = i % 2, i // 2
        left = Inches(0.6) + Inches(col * 6.3)
        top = Inches(1.4) + Inches(row * 1.6)
        card(s, left, top, Inches(6.0), Inches(1.35))
        accent_bar_left(s, left, top, Inches(1.35))
        add_textbox(s, left + Inches(0.3), top + Inches(0.3), Inches(5.4), Inches(0.4), name, 15, True, NAVY, font="Consolas")
        add_textbox(s, left + Inches(0.3), top + Inches(0.75), Inches(5.4), Inches(0.35), desc, 13, False, MUTED)
    footer(s, 21)

    # 22 SQL highlights
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "SQL Highlights (Latest Run)")
    metrics = [
        ("38.4%", "High traffic delay rate"),
        ("$65,182", "Taxi total revenue (leader)"),
        ("$35.47", "Airport avg fare (premium)"),
    ]
    for i, (val, label) in enumerate(metrics):
        left = Inches(0.7) + Inches(i * 4.15)
        card(s, left, Inches(1.5), Inches(3.9), Inches(2.4))
        add_textbox(s, left, Inches(1.9), Inches(3.9), Inches(0.7), val, 32, True, TEAL, PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.2), Inches(2.8), Inches(3.5), Inches(0.7), label, 14, False, SLATE, PP_ALIGN.CENTER)
    bullet_block(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(2.2), [
        "Airport and Downtown command higher average fares",
        "High traffic is the clearest delay hotspot for operations",
        "Taxi leads revenue; Metro/Bus compete on volume and lower fares",
    ], 16)
    footer(s, 22)

    # 23 Part 5
    section_slide(prs, "PART 5", "Machine Learning", "Predict fares, delays, and segments", 23)

    # 24 Three models
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Three-Model ML Stack")
    models = [
        ("Fare Regression", ["RandomForestRegressor", "200 trees, seed=42", "Target: total_fare", "R² = 0.981", "MAE = $1.91"]),
        ("Delay Classification", ["RandomForestClassifier", "200 trees, seed=42", "Target: is_delayed", "Accuracy = 77.9%", "Base delay rate = 30.7%"]),
        ("Trip Clustering", ["KMeans (k=4)", "Scaled features", "distance, duration,", "fare, delay_minutes", "Silhouette = 0.323"]),
    ]
    for i, (title, lines) in enumerate(models):
        left = Inches(0.55) + Inches(i * 4.2)
        card(s, left, Inches(1.4), Inches(3.95), Inches(5.15))
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.4), Inches(3.95), Inches(0.85))
        rgb_fill(hdr, NAVY if i != 1 else TEAL)
        no_line(hdr)
        add_textbox(s, left, Inches(1.6), Inches(3.95), Inches(0.5), title, 16, True, WHITE, PP_ALIGN.CENTER)
        for j, line in enumerate(lines):
            color = TEAL if j >= 3 else SLATE
            bold = j >= 3
            add_textbox(s, left + Inches(0.25), Inches(2.55) + Inches(j * 0.6), Inches(3.45), Inches(0.5), line, 14, bold, color, PP_ALIGN.CENTER)
    footer(s, 24)

    # 25 Preprocessing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Preprocessing & Leakage Control")
    bullet_block(s, Inches(0.7), Inches(1.5), Inches(11.5), Inches(5), [
        "ColumnTransformer: StandardScaler on numeric features",
        "OneHotEncoder(handle_unknown='ignore') on categoricals",
        "Features: hour, vehicle, zone, distance, duration, passengers, payment, weather, traffic",
        "Excluded from X: total_fare, is_delayed, tip_amount, delay_minutes, trip_id, trip_date",
        "Pipeline = preprocessor + model; train_test_split 80/20, random_state=42",
        "Same split seed for regression and classification → comparable evaluation",
    ], 16)
    footer(s, 25)

    # 26 Metrics interpretation
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Interpreting Model Performance")
    cards_m = [
        ("R² 0.981", "Distance, vehicle type, and zone dominate fare structure"),
        ("MAE $1.91", "Predictions are close enough for pricing analytics demos"),
        ("Accuracy 77.9%", "Useful ops signal despite stochastic traffic/weather"),
        ("Silhouette 0.323", "Moderate separation: short urban vs longer airport patterns"),
    ]
    for i, (t, d) in enumerate(cards_m):
        col, row = i % 2, i // 2
        left = Inches(0.6) + Inches(col * 6.3)
        top = Inches(1.4) + Inches(row * 2.4)
        card(s, left, top, Inches(6.0), Inches(2.15))
        accent_bar_left(s, left, top, Inches(2.15))
        add_textbox(s, left + Inches(0.35), top + Inches(0.4), Inches(5.3), Inches(0.5), t, 20, True, TEAL)
        add_textbox(s, left + Inches(0.35), top + Inches(1.1), Inches(5.3), Inches(0.7), d, 14, False, SLATE)
    footer(s, 26)

    # 27 Visuals fare
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Visual Insights — Fare & Distance")
    fare = FIG / "fare_distribution.png"
    dist = FIG / "distance_vs_fare.png"
    if fare.exists():
        s.shapes.add_picture(str(fare), Inches(0.5), Inches(1.35), height=Inches(5.3))
    if dist.exists():
        s.shapes.add_picture(str(dist), Inches(6.7), Inches(1.35), height=Inches(5.3))
    footer(s, 27)

    # 28 Visuals delay/revenue
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Visual Insights — Delays & Revenue")
    delay = FIG / "delay_by_traffic.png"
    rev = FIG / "revenue_by_vehicle.png"
    peak = FIG / "peak_hour_trips.png"
    if delay.exists():
        s.shapes.add_picture(str(delay), Inches(0.4), Inches(1.35), width=Inches(4.1))
    if rev.exists():
        s.shapes.add_picture(str(rev), Inches(4.6), Inches(1.35), width=Inches(4.1))
    if peak.exists():
        s.shapes.add_picture(str(peak), Inches(8.8), Inches(1.35), width=Inches(4.1))
    footer(s, 28)

    # 29 Part 6
    section_slide(prs, "PART 6", "Dashboard & Demo", "Interactive exploration for stakeholders", 29)

    # 30 Dashboard
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Streamlit Interactive Dashboard")
    left_items = [
        "KPIs: trips, fare, distance, delay, revenue",
        "ML metric cards from model_metrics.json",
        "Sidebar filters update all charts live",
        "Plotly dual-axis & scatter charts",
        "Sortable trip records table",
        "Run: py -m streamlit run dashboard.py",
        "URL: http://localhost:8501",
    ]
    bullet_block(s, Inches(0.6), Inches(1.4), Inches(6.2), Inches(5.2), left_items, 15)
    shot = PHOTO / "Screenshot 2026-06-24 225934.png"
    if shot.exists():
        s.shapes.add_picture(str(shot), Inches(7.0), Inches(1.5), width=Inches(5.7))
    footer(s, 30)

    # 31 How to run
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "How to Run the Project")
    card(s, Inches(0.55), Inches(1.45), Inches(6.0), Inches(5.1))
    add_textbox(s, Inches(0.9), Inches(1.75), Inches(5.3), Inches(0.45), "Offline Pipeline", 18, True, NAVY)
    add_textbox(s, Inches(0.9), Inches(2.5), Inches(5.3), Inches(1.5),
                "py -m pip install -r requirements.txt\n\npy src/main.py", 15, False, SLATE, font="Consolas")
    add_textbox(s, Inches(0.9), Inches(4.4), Inches(5.3), Inches(1.5),
                "Produces Parquet, SQLite,\nmetrics, figures, Markdown report", 14, False, MUTED)
    card(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(5.1))
    add_textbox(s, Inches(7.15), Inches(1.75), Inches(5.3), Inches(0.45), "Interactive Website", 18, True, NAVY)
    add_textbox(s, Inches(7.15), Inches(2.5), Inches(5.3), Inches(1.5),
                "py -m streamlit run dashboard.py\n\nor double-click run_website.bat", 15, False, SLATE, font="Consolas")
    add_textbox(s, Inches(7.15), Inches(4.4), Inches(5.3), Inches(1.5),
                "Opens localhost:8501\nFilters · KPIs · Charts · Table", 14, False, MUTED)
    footer(s, 31)

    # 32 Insights
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    content_title(s, "Key Insights & Impact")
    insights = [
        "Distance and vehicle type strongly drive fare levels",
        "Airport / Downtown are premium zones; Taxi leads revenue",
        "High traffic and adverse weather raise delay risk — actionable for ops",
        "Trip segments separate short urban hops from longer airport routes",
        "One reproducible pipeline from raw CSV to dashboard-ready insights",
        "Ready for demos, grading, and extension to real operator data",
    ]
    for i, text in enumerate(insights):
        y = Inches(1.35) + Inches(i * 0.85)
        card(s, Inches(0.7), y, Inches(11.9), Inches(0.72))
        add_textbox(s, Inches(1.0), y + Inches(0.18), Inches(11.2), Inches(0.45), f"{i+1}.  {text}", 15, False, SLATE)
    footer(s, 32)

    # 33 Thank you
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_slide(s)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), SLIDE_H)
    rgb_fill(strip, TEAL)
    no_line(strip)
    soft = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), 0, Inches(0.08), SLIDE_H)
    rgb_fill(soft, ACCENT_SOFT)
    no_line(soft)
    add_textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.9), "Thank You", 48, True, NAVY)
    add_textbox(s, Inches(0.9), Inches(3.2), Inches(11), Inches(0.5), "Questions & Discussion", 22, False, TEAL)
    add_textbox(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.6),
                "Ahmed El-Gamal  ·  Youssef Mustafa  ·  Maryam Khaled  ·  Mariam Tarek  ·  Mostafa Fouad  ·  Nada Ahmed",
                13, False, SLATE)
    add_textbox(s, Inches(0.9), Inches(5.0), Inches(11), Inches(0.4),
                "DEPI — AI & Data Science  ·  Transport Analysis Project", 14, False, MUTED)
    footer(s, 33)
    notes(s, "Invite questions. Offer a short live demo of the dashboard if time allows.")

    prs.save(str(OUT_TMP))
    try:
        if OUT.exists():
            OUT.unlink()
        OUT_TMP.rename(OUT)
        saved = OUT
    except OSError:
        saved = OUT_TMP
    print(f"Saved: {saved}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
